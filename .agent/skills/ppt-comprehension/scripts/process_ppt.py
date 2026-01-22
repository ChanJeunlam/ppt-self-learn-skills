import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_content(pptx_path, output_dir):
    """
    Extracts text and images from a PPTX file.
    
    Structure:
    output_dir/
      draft.md
      images/
        slide_1_image_1.jpg
        slide_1_image_2.png
        ...
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    images_dir = os.path.join(output_dir, "images")
    if not os.path.exists(images_dir):
        os.makedirs(images_dir)

    prs = Presentation(pptx_path)
    
    md_content = f"# Study Guide: {os.path.basename(pptx_path)}\n\n"
    md_content += "> [!NOTE]\n"
    md_content += "> Extracted content from presentation. Images are embedded below.\n\n"

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        md_content += "---\n\n"
        md_content += f"## Slide {slide_num}\n\n"
        
        slide_text = []
        slide_images = []

        # Iterate through all shapes in the slide
        # We sort shapes by potential reading order (top-to-bottom, left-to-right)
        # This is heuristics-based since PPTX shapes are Z-order based.
        
        # A simple sort key: (top, left)
        # Note: shape.top and shape.left are EMU (English Metric Units) or None
        def sort_key(s):
            top = getattr(s, 'top', 0) or 0
            left = getattr(s, 'left', 0) or 0
            return (top, left)

        sorted_shapes = sorted(slide.shapes, key=sort_key)

        for shape_idx, shape in enumerate(sorted_shapes):
            # 1. Extract Text
            if shape.has_text_frame:
                # We need to preserve paragraph structure and bullets
                shape_text_block = []
                for paragraph in shape.text_frame.paragraphs:
                    p_text = paragraph.text.strip()
                    if p_text:
                        # Check for bullet level (0 is main bullet, 1 is sub-bullet, etc.)
                        # We use a markdown list if it has a bullet char or if indentation > 0
                        # But simplest is to just behave like a list if it's not a title
                        # Let's rely on paragraph.level if available
                        level = paragraph.level
                        indent = "  " * level
                        
                        # Heuristic: If it's a short line, treat as bullet.
                        prefix = "- "
                        
                        # PPT titles usually have level 0 but are distinct. 
                        # We can just format everything as list items for "Original Content" to capture structure
                        shape_text_block.append(f"{indent}{prefix}{p_text}")
                
                if shape_text_block:
                    slide_text.extend(shape_text_block)
            
            # 2. Extract Images (Picture)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    # Get image extension
                    ext = image.ext
                    image_filename = f"slide_{slide_num}_img_{shape_idx}.{ext}"
                    image_path = os.path.join(images_dir, image_filename)
                    
                    with open(image_path, "wb") as f:
                        f.write(image.blob)
                        
                    slide_images.append({
                        "path": f"images/{image_filename}",
                        "alt": f"Slide {slide_num} Image {shape_idx}"
                    })
                except Exception as e:
                    print(f"Warning: Failed to extract image from slide {slide_num}: {e}")

            # 3. Extract Images inside Shapes (BlipFill)
            # Some shapes have picture fills.
            # TODO: Add advanced handling for shape fills if needed.

        # --- Construct Markdown for this Slide ---
        
        # A. Display Images First (Visual Context)
        if slide_images:
            md_content += "### Visuals\n\n"
            for img in slide_images:
                md_content += f"![{img['alt']}]({img['path']})\n\n"
        
        # B. Display Text
        md_content += "### Original Content\n\n"
        if slide_text:
            # Check if it looks like a list (which we formatted as such)
            # We just join them with newlines.
            for line in slide_text:
                md_content += f"{line}\n"
            md_content += "\n"
        else:
            md_content += "*[No text detected]*\n\n"
            
        # C. Explanation Placeholder
        # Unique placeholder for "In-Place Edit" workflow
        placeholder_text = f"AI: Explain Slide {slide_num} - Please analyze the text and visuals above to explain the key concepts of this slide."
        md_content += "### 🎓 讲师讲解\n\n"
        md_content += f"> [{placeholder_text}]\n\n"

    draft_path = os.path.join(output_dir, "draft.md")
    with open(draft_path, "w", encoding="utf-8") as f:
        f.write(md_content)
        
    print(f"Success! Draft saved to: {draft_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: uv run process_ppt.py <input_pptx> <output_dir>")
        sys.exit(1)
        
    pptx_file = sys.argv[1]
    out_dir = sys.argv[2]
    
    if not os.path.exists(pptx_file):
        print(f"Error: File '{pptx_file}' not found.")
        sys.exit(1)
        
    extract_content(pptx_file, out_dir)
